# -*- coding: utf-8 -*-
"""地図PNG＋ジオコーディング済みCSVから編集可能なPPTXを組み立てる

- ピン（番号入り円）・凡例・表はすべてPowerPointのネイティブシェイプ/テーブル
  として配置するため、PowerPoint上で自由に移動・修正できる
- ピン位置は Webメルカトル座標系で緯度経度→地図画像内ピクセル→スライド上の
  インチ座標に換算して決める
"""
import csv
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from common import OUTPUT_DIR, lonlat_to_world, opening_year, world_to_pixel

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
JP_FONT = "Meiryo"

CATEGORY_COLORS = {
    "ビジネスホテル": RGBColor(0x25, 0x63, 0xEB),
    "ホステル": RGBColor(0xF5, 0x9E, 0x0B),
    "シティホテル": RGBColor(0x10, 0xB9, 0x81),
    "リゾートホテル": RGBColor(0x8B, 0x5C, 0xF6),
}
DEFAULT_COLOR = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PIN_COLOR = RGBColor(0x25, 0x63, 0xEB)  # マップページのピンは単色（青）で統一

# マップスライドのレイアウト（インチ）。地図枠は 8.0in×6.2in 固定（render_mapのアスペクトと一致）
MAP_LEFT, MAP_TOP, MAP_H_IN = 0.25, 1.0, 6.2


def cat_color(category):
    return CATEGORY_COLORS.get((category or "").strip(), DEFAULT_COLOR)


def style_run(run, size, bold=False, color=DARK, name=JP_FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 日本語（East Asian）フォントはXML直接指定が必要
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def add_textbox(slide, left, top, width, height, text, size, bold=False,
                color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        style_run(p.add_run(), size, bold, color)
        p.runs[0].text = line
    return box


def load_rows(geocoded_csv):
    with open(geocoded_csv, encoding="utf-8-sig") as f:
        return [r for r in csv.DictReader(f)]


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------- スライド1: タイトル ----------------

def add_title_slide(prs, market, date_str, n_hotels):
    slide = blank_slide(prs)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), SLIDE_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = cat_color("ビジネスホテル")
    bar.line.fill.background()
    add_textbox(slide, 1.0, 2.6, 11.3, 1.0, f"{market}　ベンチマークホテル マップ", 40, True)
    add_textbox(slide, 1.0, 3.7, 11.3, 0.5,
                f"対象マーケットにおけるベンチマークホテルの分布と概要（{n_hotels}施設）", 18, False, GRAY)
    add_textbox(slide, 1.0, 6.6, 11.3, 0.7,
                f"作成日: {date_str}\n地図: © OpenStreetMap contributors ／ 位置情報: 国土地理院ジオコーディングAPI",
                10, False, GRAY)
    return slide


# ---------------- スライド2: マップ ----------------

def latlon_to_slide_inches(lat, lon, meta, map_w_in, map_h_in):
    wx, wy = lonlat_to_world(lon, lat)
    px, py = world_to_pixel(wx, wy, meta["zoom"])
    x_img = px - (meta["center_px_x"] - meta["width_px"] / 2)
    y_img = py - (meta["center_px_y"] - meta["height_px"] / 2)
    return (MAP_LEFT + x_img / meta["width_px"] * map_w_in,
            MAP_TOP + y_img / meta["height_px"] * map_h_in)


def add_pin(slide, x_in, y_in, number, color, diameter=0.30):
    d = Inches(diameter)
    pin = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x_in) - int(d / 2), Inches(y_in) - int(d / 2), d, d)
    pin.fill.solid()
    pin.fill.fore_color.rgb = color
    pin.line.color.rgb = WHITE
    pin.line.width = Pt(1.25)
    pin.shadow.inherit = False
    tf = pin.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.add_run(), 10, True, WHITE)
    p.runs[0].text = str(number)


def add_map_slide(prs, rows, map_png, meta, market):
    slide = blank_slide(prs)
    add_textbox(slide, 0.25, 0.2, 10.0, 0.5, f"{market}｜ベンチマークホテル分布マップ", 22, True)
    add_textbox(slide, 0.25, 0.68, 10.0, 0.3,
                "ピンの番号は右表のNo.に対応（ピン・表はPowerPoint上で編集可能）", 10, False, GRAY)

    map_w_in = MAP_H_IN * meta["width_px"] / meta["height_px"]
    slide.shapes.add_picture(map_png, Inches(MAP_LEFT), Inches(MAP_TOP),
                             Inches(map_w_in), Inches(MAP_H_IN))

    # ピン（座標のある施設のみ）
    for i, row in enumerate(rows, 1):
        if not row.get("緯度") or not row.get("経度"):
            continue
        x_in, y_in = latlon_to_slide_inches(float(row["緯度"]), float(row["経度"]),
                                            meta, map_w_in, MAP_H_IN)
        add_pin(slide, x_in, y_in, i, PIN_COLOR)

    # 右側パネル: 凡例テーブル
    panel_left = MAP_LEFT + map_w_in + 0.25
    panel_w = 13.333 - panel_left - 0.2
    n = len(rows)
    tbl_shape = slide.shapes.add_table(n + 1, 4, Inches(panel_left), Inches(MAP_TOP),
                                       Inches(panel_w), Inches(0.28 * (n + 1)))
    tbl = tbl_shape.table
    widths = [0.5, panel_w - 0.5 - 0.85 - 0.95, 0.85, 0.95]
    for c, w in enumerate(widths):
        tbl.columns[c].width = Inches(w)
    headers = ["No", "ホテル名", "客室数", "開業年"]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        cell.margin_left = cell.margin_right = Inches(0.04)
        cell.margin_top = cell.margin_bottom = Inches(0.01)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style_run(p.add_run(), 9, True, WHITE)
        p.runs[0].text = h
    for i, row in enumerate(rows, 1):
        rooms = row.get("部屋数", "")
        values = [str(i), row.get("施設名", ""), f"{rooms}室" if rooms else "-",
                  opening_year(row.get("開業"))]
        for c, v in enumerate(values):
            cell = tbl.cell(i, c)
            cell.margin_left = cell.margin_right = Inches(0.04)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF3, 0xF4, 0xF6)
            if c == 0:  # No列はピンと同じ色（単色）で塗る
                cell.fill.fore_color.rgb = PIN_COLOR
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c != 1 else PP_ALIGN.LEFT
            style_run(p.add_run(), 8.5, c == 0, WHITE if c == 0 else DARK)
            p.runs[0].text = v
        tbl.rows[i].height = Inches(0.28)
    tbl.rows[0].height = Inches(0.28)

    # 位置特定できなかった施設の注記（テーブルの下）
    notice_top = MAP_TOP + 0.28 * (n + 1) + 0.2
    missing = [f"{i}. {r.get('施設名')}" for i, r in enumerate(rows, 1) if not r.get("緯度")]
    if missing:
        add_textbox(slide, panel_left, notice_top, panel_w, 0.8,
                    "※位置特定不可: " + "、".join(missing), 8, False, GRAY)
    return slide


# ---------------- スライド3: 詳細一覧 ----------------

def add_table_slide(prs, rows, market):
    slide = blank_slide(prs)
    add_textbox(slide, 0.25, 0.2, 12.0, 0.5, f"{market}｜ベンチマークホテル一覧", 22, True)

    # 全データにプライス・インデックスが無ければ列を省く（空列を出さない）
    has_price = any((r.get("プライス・インデックス") or "").strip() for r in rows)

    # 列定義: (見出し, 幅in, 揃え, 値関数)。price列が無い分の幅は施設名・住所に配分
    columns = [
        ("No", 0.5, PP_ALIGN.CENTER, lambda i, r: str(i)),
        ("施設名", 3.0 if has_price else 3.7, PP_ALIGN.LEFT, lambda i, r: r.get("施設名", "")),
        ("カテゴリー", 1.4, PP_ALIGN.CENTER, lambda i, r: r.get("カテゴリー", "")),
        ("部屋数", 0.8, PP_ALIGN.CENTER, lambda i, r: (r.get("部屋数") or "-")),
        ("開業年", 0.9, PP_ALIGN.CENTER, lambda i, r: opening_year(r.get("開業"))),
    ]
    if has_price:
        columns.append(("プライス・インデックス(円)", 2.0, PP_ALIGN.LEFT,
                        lambda i, r: r.get("プライス・インデックス", "")))
    columns.append(("住所", 4.2 if has_price else 4.5, PP_ALIGN.LEFT, lambda i, r: r.get("住所", "")))

    headers = [c[0] for c in columns]
    widths = [c[1] for c in columns]
    n = len(rows)
    total_w = sum(widths)
    tbl_shape = slide.shapes.add_table(n + 1, len(columns), Inches(0.25), Inches(0.95),
                                       Inches(total_w), Inches(0.32 * (n + 1)))
    tbl = tbl_shape.table
    for c, w in enumerate(widths):
        tbl.columns[c].width = Inches(w)
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style_run(p.add_run(), 9.5, True, WHITE)
        p.runs[0].text = h
    for i, row in enumerate(rows, 1):
        for c, (_, _, align, fn) in enumerate(columns):
            cell = tbl.cell(i, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF3, 0xF4, 0xF6)
            p = cell.text_frame.paragraphs[0]
            p.alignment = align
            style_run(p.add_run(), 9, False, DARK)
            p.runs[0].text = fn(i, row)
        tbl.rows[i].height = Inches(0.32)
    note = "出典: 各施設公表情報等"
    if has_price:
        note += " ／ プライス・インデックスは調査時点の参考価格帯"
    add_textbox(slide, 0.25, 7.05, 12.8, 0.3, note, 8.5, False, GRAY)
    return slide


def main(geocoded_csv, map_png, map_meta, out_pptx, market, date_str):
    rows = load_rows(geocoded_csv)
    meta = json.loads(Path(map_meta).read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    add_title_slide(prs, market, date_str, len(rows))
    add_map_slide(prs, rows, map_png, meta, market)
    add_table_slide(prs, rows, market)
    prs.save(out_pptx)
    print(f"出力: {out_pptx}（{len(rows)}施設）")


if __name__ == "__main__":
    import datetime
    args = sys.argv[1:]
    geocoded = args[0] if len(args) > 0 else str(OUTPUT_DIR / "geocoded.csv")
    png = args[1] if len(args) > 1 else str(OUTPUT_DIR / "map.png")
    meta = args[2] if len(args) > 2 else str(OUTPUT_DIR / "map_meta.json")
    out = args[3] if len(args) > 3 else str(OUTPUT_DIR / "benchmark_hotels.pptx")
    market = args[4] if len(args) > 4 else "那覇市"
    main(geocoded, png, meta, out, market, datetime.date.today().isoformat())
